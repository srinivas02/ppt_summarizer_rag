import io
import os
import traceback
import pandas as pd

from tqdm import tqdm
from PIL import Image
from pptx.enum.chart import *
from pptx.enum.shapes import *

from pptx import Presentation

from haystack.schema import Document as hay_doc
from langchain.schema.document import Document as lc_doc
from mavericks.common.utils import chartTypeGlossary


class PPTXExtractor:
    """
    Extracts content from a PowerPoint (PPTX) file, including titles, text, charts, tables, and images.
    """

    def __init__(self, filepath):
        """
        Initializes the PPTXExtractorFixed instance with the given file path.

        Parameters:
        - filepath (str): The path to the PowerPoint file.
        """
        self.filepath = filepath
        self.presentation = Presentation(self.filepath)
        self.metadata = {
                "source": self.filepath,
                'file_directory': os.path.dirname(self.filepath),
                'file_name': os.path.basename(self.filepath),
                'last_modified': os.path.getmtime(self.filepath)
            }

    def get_slide_content(self, return_langchain_documents=True):
        """
        Extracts content from each slide in the presentation.

        Returns:
        - list: List of dictionaries containing slide content.
        """
        slide_contents = []
        for slide_number, slide in tqdm(enumerate(self.presentation.slides, start=1)):
            slide_content = {
                "Slide #": slide_number,
                "Layout": self.extract_layout(slide),
                "Title": self.extract_title(slide),
                "Text Content": self.extract_text(slide),
                "Charts": self.extract_charts(slide),
                "Tables": self.extract_tables(slide)
            }
            charts = self.extract_charts(slide)
            tables = self.extract_tables(slide)
            # metadata=datapoint.metadata,
            if tables is not None and charts is not None:
                datasource = f"Data Source:\n\nCharts Data:\n{charts}\n\nTables Data:\n{tables}"
            elif charts is not None:
                datasource = f"Data Source:\n\nCharts Data:\n{charts}"
            elif tables is not None:
                datasource = f"Data Source:\n\nTables Data:\n{tables}"
            else:
                datasource = ""
            content=f"""Slide #:{slide_number}\n\n\nTitle: {str(self.extract_title(slide))}\n\n\nContent:{str(self.extract_text(slide))}\n\n\n{str(datasource)}"""
            if return_langchain_documents:
              document_data = lc_doc(page_content=content, metadata=self.metadata)
            else:
              document_data = hay_doc(content=content, meta=self.metadata)
            slide_contents.append(document_data)
        return slide_contents

    def extract_layout(self, slide):
        """
        Extracts the layout details from the given slide.

        Parameters:
        - slide: PowerPoint slide object.

        Returns:
        - str: layout name.
        """
        return slide.slide_layout.name

    def extract_title(self, slide):
        """
        Extracts the title from the given slide.

        Parameters:
        - slide: PowerPoint slide object.

        Returns:
        - str: Title text or 'No Title' if no title is found.
        """
        title = slide.shapes.title
        return title.text if title else 'No Title'

    def extract_text(self, slide):
        """
        Extracts text content from the given slide.

        Parameters:
        - slide: PowerPoint slide object.

        Returns:
        - str: Concatenated text content.
        """
        text_content = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text_content.append(paragraph.text)
        return "\n".join(text_content)

    def extract_charts(self, slide):
        """
        Extracts chart details from the given slide.

        Parameters:
        - slide: PowerPoint slide object.

        Returns:
        - list: List of dictionaries containing chart details.
        """
        chart_details = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = shape.chart
                chart_type = next((key for key, values in chartTypeGlossary.items() if
                                   str(chart.chart_type).split(" ")[0] in values), "Unknown Chart Type")
                chart_data = self.extract_chart_data(chart)
                chart_details.append({
                    "Chart Type": chart_type,
                    "Chart Data": chart_data
                })
        return chart_details if chart_details else 'No Chart'

    def extract_chart_data(self, chart) -> pd.DataFrame:
        """
        Extracts data from the embedded Excel workbook of the given chart.

        Parameters:
        - chart: Chart object.

        Returns:
        - pd.DataFrame or str: DataFrame containing chart data or 'No Data' if no data is found.
        """
        if chart.part.chart_workbook:
            xlsx_blob = chart.part.chart_workbook.xlsx_part.blob
            df = pd.read_excel(io.BytesIO(xlsx_blob),
                               sheet_name=0, header=None)
            return df
        return 'No Data'

    def extract_tables(self, slide):
        """
        Extracts table details from the given slide.

        Parameters:
        - slide: PowerPoint slide object.

        Returns:
        - list: List of DataFrames containing table details.
        """
        table_details = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_data = [[cell.text.strip() for cell in row.cells]
                              for row in table.rows]
                df = pd.DataFrame(table_data)
                table_details.append(df)
        return table_details if table_details else 'No Table'

    

# # Example usage:
# pptx_path = "/content/drive/MyDrive/ppts/DTC_Trends.pptx"
# extractor_fixed = PPTXExtractorFixed(pptx_path)
# slide_contents_fixed = extractor_fixed.get_slide_content()
